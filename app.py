from datetime import datetime
import io
import os
import psycopg2


from deep_translator import GoogleTranslator, MyMemoryTranslator
from pypdf import PdfReader
from pptx import Presentation
from PIL import Image
import streamlit as st

try:
    import pytesseract
    _TESSERACT_IMPORTED = True
except ImportError:
    _TESSERACT_IMPORTED = False

try:
    import speech_recognition as sr
    _SPEECH_RECOGNITION_IMPORTED = True
except ImportError:
    _SPEECH_RECOGNITION_IMPORTED = False

try:
    from gtts import gTTS
    _GTTS_IMPORTED = True
except ImportError:
    _GTTS_IMPORTED = False

# --- 1. DATABASE SETUP & HELPERS ---

DATABASE_URL = os.getenv("DATABASE_URL")

def get_db_connection():
    return psycopg2.connect(DATABASE_URL)

def init_db():
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS translations (
            id SERIAL PRIMARY KEY,
            user_role TEXT NOT NULL,
            field TEXT NOT NULL,
            source_text TEXT NOT NULL,
            translated_text TEXT NOT NULL,
            target_language TEXT NOT NULL,
            timestamp TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()

init_db()

def save_to_db(user_role, field, source_text, translated_text, target_language):
    conn = get_db_connection()
    c = conn.cursor()
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    c.execute("""
        INSERT INTO translations 
            (user_role, field, source_text, translated_text, target_language, timestamp)
        VALUES (%s, %s, %s, %s, %s, %s)
    """, (user_role, field, source_text, translated_text, target_language, timestamp))
    conn.commit()
    conn.close()

def get_all_records():
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT * FROM translations ORDER BY id DESC")
    rows = c.fetchall()
    conn.close()
    return rows

def delete_record(record_id):
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("DELETE FROM translations WHERE id = %s", (record_id,))
    conn.commit()
    conn.close()

        

 

  
     # --- 1. DATABASE SETUP & HELPERS ---

DATABASE_URL = os.getenv("DATABASE_URL")

def get_db_connection():
    return psycopg2.connect(DATABASE_URL)

def init_db():
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS translations (
            id SERIAL PRIMARY KEY,
            user_role TEXT NOT NULL,
            field TEXT NOT NULL,
            source_text TEXT NOT NULL,
            translated_text TEXT NOT NULL,
            target_language TEXT NOT NULL,
            timestamp TEXT NOT NULL
        )
    """)
    conn.commit()
    conn.close()

init_db()

def save_to_db(user_role, field, source_text, translated_text, target_language):
    conn = get_db_connection()
    c = conn.cursor()
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    c.execute("""
        INSERT INTO translations 
            (user_role, field, source_text, translated_text, target_language, timestamp)
        VALUES (%s, %s, %s, %s, %s, %s)
    """, (user_role, field, source_text, translated_text, target_language, timestamp))
    conn.commit()
    conn.close()

def get_all_records():
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("SELECT * FROM translations ORDER BY id DESC")
    rows = c.fetchall()
    conn.close()
    return rows

def delete_record(record_id):
    conn = get_db_connection()
    c = conn.cursor()
    c.execute("DELETE FROM translations WHERE id = %s", (record_id,))
    conn.commit()
    conn.close()

    conn.commit()
    conn.close()


# --- FILE TEXT EXTRACTION HELPERS ---

def extract_text_from_pdf(file_bytes):
    """Pull all readable text out of a PDF, page by page. Returns None on failure."""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        pages_text = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages_text.append(page_text.strip())
        return "\n\n".join(pages_text).strip() or None
    except Exception:
        return None


def extract_text_from_pptx(file_bytes):
    """Pull all text boxes/bullets out of a PowerPoint deck, slide by slide. Returns None on failure."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            lines.append(line)
            if lines:
                slides_text.append(f"Slide {i}:\n" + "\n".join(lines))
        return "\n\n".join(slides_text).strip() or None
    except Exception:
        return None


def extract_text_from_image(file_bytes):
    """
    OCR an image (e.g. a photo of handwritten/printed notes) using
    Tesseract via pytesseract. Returns None if OCR isn't available on
    this machine rather than raising, so the caller can show a helpful
    message instead of crashing.
    """
    if not _TESSERACT_IMPORTED:
        return None
    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception:
        return None


# Language Code Mapper for Google Translate
LANG_CODES = {
    "isiZulu": "zu",
    "isiXhosa": "xh",
    "Afrikaans": "af",
    "English": "en",
    "Sepedi": "nso",
    "Setswana": "tn",
    "Sesotho": "st",
    "Xitsonga": "ts",
    "siSwati": "ss",
    "Tshivenda": "ve",
    "isiNdebele": "nr",
    "South African Sign Language (SASL)": "en",  # fallback display mapping
}

# MyMemory uses full locale-style codes and only covers a subset of
# South Africa's official languages. Languages not listed here have no
# fallback if Google is unavailable.
MYMEMORY_CODES = {
    "isiZulu": "zu-ZA",
    "isiXhosa": "xh-ZA",
    "Afrikaans": "af-ZA",
    "English": "en-GB",
    "Setswana": "tn-BW",
    "Sesotho": "st-ST",
}


def _looks_like_error_page(text):
    """Detect an HTML/server error page returned as if it were a translation."""
    if not text:
        return True
    lowered = text.lower()
    return "<html" in lowered or "server error" in lowered or "error 500" in lowered


def _looks_untranslated(original, translated, target_lang_name):
    """
    Detect a translation call that silently echoed the input back instead
    of translating it - a known failure mode of free translation
    endpoints when text is too long or the service is having trouble.
    Not a perfect check (e.g. a single shared proper noun could match),
    but effective for whole-block comparisons.
    """
    if target_lang_name == "English":
        return False  # translating into English can legitimately be a no-op
    return translated.strip().lower() == original.strip().lower()


def _split_into_chunks(text, max_len):
    """
    Split text into pieces no longer than max_len, preferring to break on
    line boundaries so slide/paragraph structure survives. Falls back to
    hard character splitting for any single line longer than max_len.
    """
    chunks = []
    current_lines = []
    current_len = 0

    def flush():
        if current_lines:
            chunks.append("\n".join(current_lines))

    for line in text.split("\n"):
        if len(line) > max_len:
            flush()
            current_lines.clear()
            current_len = 0
            for i in range(0, len(line), max_len):
                chunks.append(line[i : i + max_len])
            continue

        added_len = len(line) + 1
        if current_len + added_len > max_len and current_lines:
            flush()
            current_lines = [line]
            current_len = added_len
        else:
            current_lines.append(line)
            current_len += added_len

    flush()
    return [c for c in chunks if c.strip()]


# Conservative request-size limits for each free backend.
GOOGLE_CHUNK_LIMIT = 4500
MYMEMORY_CHUNK_LIMIT = 480


# --- SPEECH-TO-TEXT & TEXT-TO-SPEECH HELPERS ---

# Locale codes for Google's free Web Speech recognition backend. Coverage of
# South Africa's official languages varies and isn't officially documented,
# so recognition for the less common languages below is best-effort - the
# caller falls back to a friendly error message if a given locale isn't
# actually supported by the backend at request time.
STT_LANG_CODES = {
    "isiZulu": "zu-ZA",
    "isiXhosa": "xh-ZA",
    "Afrikaans": "af-ZA",
    "English": "en-ZA",
    "Sepedi": "nso-ZA",
    "Setswana": "tn-ZA",
    "Sesotho": "st-ZA",
    "Xitsonga": "ts-ZA",
    "siSwati": "ss-ZA",
    "Tshivenda": "ve-ZA",
    "isiNdebele": "nr-ZA",
}
SPOKEN_LANGUAGE_OPTIONS = list(STT_LANG_CODES.keys())

# gTTS (Google Text-to-Speech) only reliably supports a small subset of
# South Africa's official languages today. Languages not listed here have
# no audio playback option yet.
TTS_LANG_CODES = {
    "Afrikaans": "af",
    "English": "en",
}


def speech_to_text(audio_bytes, spoken_language):
    """
    Transcribe recorded audio (WAV bytes, e.g. from st.audio_input) to text
    using SpeechRecognition's free Google Web Speech API backend. Returns a
    (text, error_message) tuple - exactly one of which is None/empty.
    """
    if not _SPEECH_RECOGNITION_IMPORTED:
        return None, (
            "Speech-to-text needs the `SpeechRecognition` Python package. "
            "Install it with `pip install SpeechRecognition` and restart the app."
        )

    locale_code = STT_LANG_CODES.get(spoken_language, "en-ZA")
    recognizer = sr.Recognizer()
    try:
        with sr.AudioFile(io.BytesIO(audio_bytes)) as source:
            audio_data = recognizer.record(source)
        text = recognizer.recognize_google(audio_data, language=locale_code)
        return text.strip(), None
    except sr.UnknownValueError:
        return None, "Couldn't make out any speech in that recording. Please try again, speaking clearly."
    except sr.RequestError as e:
        return None, f"Speech recognition service is unavailable right now ({e}). Please try again shortly."
    except Exception as e:
        return None, f"Couldn't transcribe that recording ({e}). Please try again or type the text manually."


def text_to_speech(text, spoken_language):
    """
    Convert text to spoken audio (MP3 bytes) using gTTS. Returns an
    (audio_bytes, error_message) tuple - exactly one of which is None.
    """
    if not text or not text.strip():
        return None, "There's no text to read aloud yet."
    if not _GTTS_IMPORTED:
        return None, (
            "Text-to-speech needs the `gTTS` Python package. Install it "
            "with `pip install gTTS` and restart the app."
        )

    tts_code = TTS_LANG_CODES.get(spoken_language)
    if tts_code is None:
        return None, (
            f"Audio playback isn't available yet for {spoken_language}. "
            "This currently works for English and Afrikaans; support for the "
            "other official languages may be added as compatible free "
            "speech services become available."
        )

    try:
        buffer = io.BytesIO()
        gTTS(text=text, lang=tts_code).write_to_fp(buffer)
        return buffer.getvalue(), None
    except Exception as e:
        return None, f"Couldn't generate audio right now ({e}). Please try again shortly."


def _translate_with_google(text, code):
    translated = GoogleTranslator(source="auto", target=code).translate(text)
    if not translated or _looks_like_error_page(translated):
        raise ValueError("Google Translate returned an unexpected response.")
    return translated


def _translate_with_mymemory(text, mymemory_code):
    translated = MyMemoryTranslator(
        source=MYMEMORY_CODES["English"], target=mymemory_code
    ).translate(text)
    if not translated or _looks_like_error_page(translated):
        raise ValueError("MyMemory returned an unexpected response.")
    return translated


def perform_translation(text, target_lang_name):
    code = LANG_CODES.get(target_lang_name, "en")

    if target_lang_name == "South African Sign Language (SASL)":
        return f"[SASL GLOSS FORMAT]\nSigns mapped for: {text}"

    mymemory_code = MYMEMORY_CODES.get(target_lang_name)

    # Long text (e.g. extracted from a multi-slide deck or PDF) can exceed
    # what the free translation endpoints accept in one request. Chunk it
    # so each request stays under the relevant limit.
    google_chunks = _split_into_chunks(text, GOOGLE_CHUNK_LIMIT)

    translated_chunks = []
    google_error = None
    google_failed = False

    for chunk in google_chunks:
        try:
            result = _translate_with_google(chunk, code)
            if _looks_untranslated(chunk, result, target_lang_name):
                raise ValueError(
                    "Google Translate returned the text unchanged - likely "
                    "blocked, rate-limited, or the request was too large."
                )
            translated_chunks.append(result)
        except Exception as e:
            google_error = str(e)
            google_failed = True
            break

    if not google_failed:
        return "\n".join(translated_chunks)

    if mymemory_code is None:
        return (
            f"Translation error: Google Translate is currently unavailable "
            f"({google_error}), and there is no fallback translator for "
            f"{target_lang_name} yet. Please try again in a moment, or "
            "translate a shorter excerpt."
        )

    # Retry from scratch with MyMemory, chunked to its much smaller limit.
    mymemory_chunks = _split_into_chunks(text, MYMEMORY_CHUNK_LIMIT)
    translated_chunks = []
    mymemory_error = None

    for chunk in mymemory_chunks:
        try:
            result = _translate_with_mymemory(chunk, mymemory_code)
            if _looks_untranslated(chunk, result, target_lang_name):
                raise ValueError("MyMemory returned the text unchanged.")
            translated_chunks.append(result)
        except Exception as e:
            mymemory_error = str(e)
            translated_chunks = None
            break

    if translated_chunks is not None:
        return "\n".join(translated_chunks)

    return (
        "Translation error: both translation services are currently "
        f"unavailable (Google: {google_error} | MyMemory: {mymemory_error}). "
        "This can happen with very long text or a slow/blocked connection - "
        "try a shorter excerpt, or check your internet connection and try again."
    )


# --- 2. PAGE CONFIGURATION & CONSTANTS ---

st.set_page_config(page_title="STEAM App - South Africa", page_icon="S", layout="wide")

OFFICIAL_LANGUAGES = list(LANG_CODES.keys())
STEAM_FIELDS = [
    "Arts",
    "Science",
    "Technology",
    "Engineering",
    "Mathematics",
]

# --- 3. SIDEBAR NAVIGATION & ROLE SELECTION ---

st.sidebar.title("STEAM Portal Access")
user_role = st.sidebar.selectbox("Select Your Role", ["Student", "Teacher"])
menu = st.sidebar.radio(
    "Navigation Menu",
    [
        "Dashboard Overview",
        "Translation & Live Class Hub",
        "Study Vault & History",
    ],
)
st.sidebar.markdown("---")
st.sidebar.info(
    f"Currently logged in as **{user_role}**. Empowering education across all "
    "12 official South African languages."
)

# --- 4. MODULES ---

# MODULE A: Dashboard Overview
if menu == "Dashboard Overview":
    st.title(f"STEAM App Dashboard - {user_role} Portal")
    st.write(
        "Welcome to the indigenous language translation, study slide management, "
        "and live classroom streaming workspace."
    )

    if user_role == "Teacher":
        st.success(
            "**Teacher Mode Active:** You can paste your virtual class links "
            "(Zoom, Teams, Meet), manage curriculum domains, and broadcast "
            "real-time subtitles directly to students in their preferred languages."
        )
    else:
        st.info(
            "**Student Mode Active:** You can input custom words, upload lecture "
            "slide texts, pick any of the 12 official languages, and save "
            "translation logs into your personal offline study vault."
        )

    st.markdown("### Covered STEAM Fields & Domains")
    steam_field_cards = ["Arts", "Science", "Technology", "Engineering", "Mathematics"]
    field_cols = st.columns(5)
    for col, field_name in zip(field_cols, steam_field_cards):
        with col:
            st.markdown(
                f"""
                <div style="
                    text-align:center;
                    padding:14px 6px;
                    border-radius:10px;
                    background-color:#f0f2f6;
                    border:1px solid #e0e2e6;
                ">
                    <div style="font-size:0.78rem; color:#6c757d; margin-bottom:4px;">
                        Field
                    </div>
                    <div style="font-size:1.05rem; font-weight:600; line-height:1.25; color:#262730;">
                        {field_name}
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

# MODULE B: Translation & Live Class Hub
elif menu == "Translation & Live Class Hub":
    st.title("STEAM Translation & Live Session Hub")

    col_f, col_l = st.columns(2)
    with col_f:
        selected_field = st.selectbox("Select STEAM Field", STEAM_FIELDS)
    with col_l:
        target_lang = st.selectbox("Select Preferred / Target Language", OFFICIAL_LANGUAGES)

    st.markdown("---")

    if user_role == "Student":
        st.subheader("Student Custom Word & Slide Notes Translator")
        st.write(
            "Enter any vocabulary word, concept phrase, or paste text extracted "
            "from lecture slides below - or upload a file to extract the text "
            "automatically."
        )

        uploaded_file = st.file_uploader(
            "Optional: Upload a PDF, PowerPoint slide deck (.pptx), or image of your notes",
            type=["pdf", "pptx", "png", "jpg", "jpeg"],
        )

        if uploaded_file is None:
            # Uploader was cleared - drop any stale status from a previous file.
            st.session_state.pop("extraction_status", None)
            st.session_state.pop("last_uploaded_name", None)

        if uploaded_file is not None and st.session_state.get("last_uploaded_name") != uploaded_file.name:
            file_bytes = uploaded_file.getvalue()
            suffix = uploaded_file.name.lower().rsplit(".", 1)[-1]

            with st.spinner(f"Extracting text from {uploaded_file.name}..."):
                if suffix == "pdf":
                    extracted = extract_text_from_pdf(file_bytes)
                elif suffix == "pptx":
                    extracted = extract_text_from_pptx(file_bytes)
                else:
                    extracted = extract_text_from_image(file_bytes)

            st.session_state["last_uploaded_name"] = uploaded_file.name

            if extracted:
                st.session_state["source_input"] = extracted
                st.session_state["extraction_status"] = (
                    "success",
                    f"Extracted text from {uploaded_file.name}. Review or edit it "
                    "below before translating.",
                )
            elif suffix in ("png", "jpg", "jpeg") and not _TESSERACT_IMPORTED:
                st.session_state["extraction_status"] = (
                    "warning",
                    "Image text extraction (OCR) needs the Tesseract OCR engine "
                    "installed on this computer, and the `pytesseract` Python "
                    "package. Install Tesseract from "
                    "https://github.com/UB-Mannheim/tesseract/wiki (Windows), "
                    "then run `pip install pytesseract` and restart the app. "
                    "In the meantime, you can type or paste the notes manually below.",
                )
            else:
                st.session_state["extraction_status"] = (
                    "warning",
                    f"Couldn't find readable text in {uploaded_file.name}. "
                    "It may be a scanned/image-only PDF, an empty slide deck, "
                    "a password-protected/corrupted file, or too blurry to "
                    "read. Try another file or type the notes manually below.",
                )

        # Keep showing the last extraction result as long as that file is
        # still selected in the uploader, instead of only on the run it
        # happened - otherwise clicking Translate makes the message vanish
        # even though nothing about the extraction changed.
        if uploaded_file is not None and st.session_state.get("extraction_status"):
            level, message = st.session_state["extraction_status"]
            getattr(st, level)(message)

        with st.expander("Or speak instead of typing (voice input)"):
            spoken_lang_student = st.selectbox(
                "Language you'll be speaking",
                SPOKEN_LANGUAGE_OPTIONS,
                key="spoken_lang_student",
            )
            recording = st.audio_input("Record your word or notes", key="student_recording")
            if recording is not None and st.button("Transcribe Recording", key="transcribe_student"):
                with st.spinner("Transcribing..."):
                    text, error = speech_to_text(recording.getvalue(), spoken_lang_student)
                if text:
                    st.session_state["source_input"] = text
                    st.success("Transcribed! Review or edit it below before translating.")
                else:
                    st.warning(error)

        source_input = st.text_area(
            "Source Text / Study Notes:",
            placeholder="Type word or paste slide notes here, or upload a file above...",
            key="source_input",
        )

        translated_output = None
        if st.button("Translate Term / Notes", type="primary"):
            if source_input.strip():
                with st.spinner("Translating text..."):
                    translated_result = perform_translation(source_input, target_lang)
                    translated_output = (
                        f"[{target_lang.upper()} Translation | Field: {selected_field}]\n\n"
                        f"Source Content: {source_input}\n\n"
                        f"Translated Result:\n{translated_result}"
                    )
                st.session_state["last_translation"] = translated_output
                st.session_state["last_translation_text"] = translated_result
                st.session_state["last_translation_lang"] = target_lang
                st.success("Translation generated successfully.")
                st.code(translated_output, language="text")
            else:
                st.warning("Please enter text or notes to translate.")

        if st.session_state.get("last_translation_text"):
            if st.button("🔊 Listen to Translation"):
                with st.spinner("Generating audio..."):
                    audio_bytes, tts_error = text_to_speech(
                        st.session_state["last_translation_text"],
                        st.session_state["last_translation_lang"],
                    )
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3")
                else:
                    st.info(tts_error)

        if st.button("Save Translation to Study Vault"):
            saved_output = st.session_state.get("last_translation")
            if saved_output:
                save_to_db(
                    "Student",
                    selected_field,
                    source_input,
                    saved_output,
                    target_lang,
                )
                st.toast("Saved successfully to your local SQLite vault!")
            else:
                st.warning("Translate something first before saving.")

    else:  # Teacher Portal
        st.subheader("Teacher Live Class & Subtitle Broadcaster")
        st.write(
            "Paste your online class link and stream live text-based lecture "
            "subtitles to student devices."
        )
        class_link = st.text_input(
            "Online Class Link (Zoom, Microsoft Teams, Google Meet)",
            placeholder="https://zoom.us/j/example",
        )

        col_btn1, col_btn2 = st.columns(2)
        with col_btn1:
            if st.button("Start Live Class Session", type="primary"):
                if class_link.strip():
                    st.success(
                        f"Live session active. Streaming link and subtitle feed set to {target_lang}."
                    )
                else:
                    st.warning("Please enter a valid class URL first.")
        with col_btn2:
            if st.button("End Session"):
                st.info("Class session ended.")

        st.markdown("---")
        st.markdown("### Live Lecture Subtitle Feed")

        with st.expander("Or speak the sentence instead of typing (voice input)"):
            spoken_lang_teacher = st.selectbox(
                "Language you're speaking",
                SPOKEN_LANGUAGE_OPTIONS,
                key="spoken_lang_teacher",
            )
            teacher_recording = st.audio_input("Record the current sentence", key="teacher_recording")
            if teacher_recording is not None and st.button("Transcribe Recording", key="transcribe_teacher"):
                with st.spinner("Transcribing..."):
                    text, error = speech_to_text(teacher_recording.getvalue(), spoken_lang_teacher)
                if text:
                    st.session_state["speech_input"] = text
                    st.success("Transcribed! Review it below, then broadcast.")
                else:
                    st.warning(error)

        speech_input = st.text_area(
            "Type current spoken sentence or lecture excerpt:",
            placeholder="Type sentences here to broadcast real-time translated subtitles...",
            key="speech_input",
        )

        if st.button(f"Broadcast Subtitle in {target_lang}"):
            if speech_input.strip():
                with st.spinner("Translating and broadcasting..."):
                    sub_translated = perform_translation(speech_input, target_lang)
                    sub_output = f"[LIVE SUBTITLE - {target_lang.upper()}] {sub_translated}"
                    save_to_db(
                        "Teacher",
                        selected_field,
                        speech_input,
                        sub_output,
                        target_lang,
                    )
                st.session_state["last_subtitle_text"] = sub_translated
                st.session_state["last_subtitle_lang"] = target_lang
                st.success("Live subtitle broadcasted and synced to student vaults!")
                st.code(sub_output, language="text")
            else:
                st.warning("Please type a phrase to broadcast.")

        if st.session_state.get("last_subtitle_text"):
            if st.button("🔊 Listen to Broadcasted Subtitle"):
                with st.spinner("Generating audio..."):
                    audio_bytes, tts_error = text_to_speech(
                        st.session_state["last_subtitle_text"],
                        st.session_state["last_subtitle_lang"],
                    )
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3")
                else:
                    st.info(tts_error)

# MODULE C: Study Vault & History
elif menu == "Study Vault & History":
    st.title("Offline Study Vault & Saved Records")
    st.write(
        "Access all your archived slide notes, custom word lookups, and "
        "broadcasted class transcripts locally."
    )

    records = get_all_records()

    if not records:
        st.info(
            "Your vault is currently empty. Start translating notes or broadcasting "
            "sessions to save items here."
        )
    else:
        for row in records:
            record_id, role, field, src, translated, lang, timestamp = row
            with st.expander(f"[{field}] Target: {lang} | By: {role} ({timestamp})"):
                st.write(f"**Source Text / Input:** {src}")
                st.markdown(f"**Stored Translation / Subtitle:**\n```text\n{translated}\n```")
                if st.button("Delete Record", key=f"del_{record_id}"):
                    delete_record(record_id)
                    st.success("Record deleted from vault!")
                    st.rerun()
